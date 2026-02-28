import logging
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

class LayoutManager:
    """Applies template styling (fonts, colors, layouts) to generated slides"""
    
    def __init__(self, template_presentation: Presentation):
        """
        Initialize with template presentation
        
        Args:
            template_presentation: Loaded PPT template
        """
        self.template = template_presentation
        self.default_font = "Calibri"
        self.default_font_size = Pt(18)
        self.title_font_size = Pt(32)
        self.code_font = "Courier New"
        self.code_font_size = Pt(12)
        
        # Extract template styling
        self._extract_template_defaults()
    
    def _extract_template_defaults(self):
        """Extract default styling from template slides"""
        try:
            if len(self.template.slides) > 0:
                first_slide = self.template.slides[0]
                
                # Extract font info from first slide
                for shape in first_slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    self.default_font = run.font.name
                                if run.font.size:
                                    self.default_font_size = run.font.size
                                break
                            break
                        break
            
            logger.info(f"✅ Template defaults extracted:")
            logger.info(f"   Font: {self.default_font}")
            logger.info(f"   Font Size: {self.default_font_size}")
        
        except Exception as e:
            logger.warning(f"⚠️ Could not extract template defaults: {str(e)}")
    
    def get_content_layout(self):
        """
        Get appropriate layout for content slides (title + content)
        
        Returns:
            Slide layout object
        """
        try:
            # Try to find "Title and Content" or similar layout
            for layout in self.template.slide_layouts:
                layout_name = layout.name.lower()
                if 'title' in layout_name and 'content' in layout_name:
                    return layout
                if 'title' in layout_name and 'text' in layout_name:
                    return layout
                if layout_name == 'title and content':
                    return layout
            
            # Fallback to layout index 1 (usually Title and Content)
            if len(self.template.slide_layouts) > 1:
                return self.template.slide_layouts[1]
            
            # Last resort - blank layout
            return self.template.slide_layouts[0]
        
        except Exception as e:
            logger.error(f"❌ Error getting content layout: {str(e)}")
            return self.template.slide_layouts[0]
    
    def get_blank_layout(self):
        """
        Get blank layout for custom slides
        
        Returns:
            Blank slide layout
        """
        try:
            # Try to find blank layout
            for layout in self.template.slide_layouts:
                if 'blank' in layout.name.lower():
                    return layout
            
            # Blank layout often at index 6
            if len(self.template.slide_layouts) > 6:
                return self.template.slide_layouts[6]
            
            # Fallback
            return self.template.slide_layouts[-1]
        
        except Exception as e:
            logger.error(f"❌ Error getting blank layout: {str(e)}")
            return self.template.slide_layouts[0]
    
    def apply_title_styling(self, text_frame, title_text: str):
        """
        Apply template styling to slide title
        
        Args:
            text_frame: Text frame to style
            title_text: Title text content
        """
        try:
            text_frame.text = title_text
            
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.bold = True
                paragraph.font.size = self.title_font_size
                paragraph.font.name = self.default_font
        
        except Exception as e:
            logger.error(f"❌ Error applying title styling: {str(e)}")
    
    def apply_bullet_styling(self, text_frame, bullets: list):
        """
        Apply template styling to bullet points
        
        Args:
            text_frame: Text frame to style
            bullets: List of bullet point strings
        """
        try:
            text_frame.clear()
            
            for bullet in bullets:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = self.default_font_size
                p.font.name = self.default_font
        
        except Exception as e:
            logger.error(f"❌ Error applying bullet styling: {str(e)}")
    
    def add_code_box(self, slide, code_snippet: str, left: float = 1.0, top: float = 4.0, 
                     width: float = 8.0, height: float = 2.5):
        """
        Add code snippet box to slide with monospace font
        
        Args:
            slide: Slide object
            code_snippet: Code text to display
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
        """
        try:
            if not code_snippet or not code_snippet.strip():
                return
            
            # Add text box for code
            left_emu = Inches(left)
            top_emu = Inches(top)
            width_emu = Inches(width)
            height_emu = Inches(height)
            
            textbox = slide.shapes.add_textbox(left_emu, top_emu, width_emu, height_emu)
            text_frame = textbox.text_frame
            text_frame.word_wrap = True
            
            # Add code text
            p = text_frame.paragraphs[0]
            p.text = code_snippet
            p.font.name = self.code_font
            p.font.size = self.code_font_size
            
            # Style the text box
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
            textbox.line.color.rgb = RGBColor(200, 200, 200)  # Gray border
            
            logger.info(f"   ✅ Added code box to slide")
        
        except Exception as e:
            logger.error(f"❌ Error adding code box: {str(e)}")
    
    def create_content_slide(self, slide_data: Dict[str, Any]) -> Any:
        """
        Create a fully styled content slide from slide data
        
        Args:
            slide_data: Dict with title, bullets, code_snippet
            
        Returns:
            Styled slide object
        """
        try:
            # Get appropriate layout
            layout = self.get_content_layout()
            slide = self.template.slides.add_slide(layout)
            
            # Find title and content placeholders
            title_shape = None
            content_shape = None
            
            for shape in slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == 1:  # Title placeholder
                        title_shape = shape
                    elif phf.type == 2:  # Content placeholder
                        content_shape = shape
            
            # Apply title
            if title_shape and 'title' in slide_data:
                self.apply_title_styling(title_shape.text_frame, slide_data['title'])
            
            # Apply bullets
            if content_shape and 'bullets' in slide_data:
                self.apply_bullet_styling(content_shape.text_frame, slide_data['bullets'])
            
            # Add code snippet if present
            if slide_data.get('code_snippet'):
                self.add_code_box(slide, slide_data['code_snippet'])
            
            logger.info(f"✅ Created slide: {slide_data.get('title', 'Untitled')}")
            
            return slide
        
        except Exception as e:
            logger.error(f"❌ Error creating content slide: {str(e)}")
            return None
    
    def create_title_slide(self, title: str, subtitle: str = "") -> Any:
        """
        Create title slide with template styling
        
        Args:
            title: Main title text
            subtitle: Subtitle text
            
        Returns:
            Title slide object
        """
        try:
            # Use first layout (usually title slide)
            title_layout = self.template.slide_layouts[0]
            slide = self.template.slides.add_slide(title_layout)
            
            # Find title and subtitle placeholders
            for shape in slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == 1:  # Title
                        shape.text = title
                    elif phf.type == 2:  # Subtitle
                        shape.text = subtitle
            
            logger.info(f"✅ Created title slide: {title}")
            
            return slide
        
        except Exception as e:
            logger.error(f"❌ Error creating title slide: {str(e)}")
            return None
    
    def apply_consistent_styling(self, slides: list):
        """
        Apply consistent styling across all slides
        
        Args:
            slides: List of slide objects
        """
        try:
            for slide in slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            # Ensure consistent font
                            if not paragraph.runs:
                                continue
                            
                            for run in paragraph.runs:
                                if not run.font.name:
                                    run.font.name = self.default_font
                                if not run.font.size:
                                    run.font.size = self.default_font_size
            
            logger.info(f"✅ Applied consistent styling to {len(slides)} slides")
        
        except Exception as e:
            logger.error(f"❌ Error applying consistent styling: {str(e)}")
