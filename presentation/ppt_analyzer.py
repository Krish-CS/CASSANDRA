import logging
from typing import Dict, List, Any
import re
from pptx import Presentation
import os

logger = logging.getLogger(__name__)

class PPTAnalyzer:
    """Analyzes PPT template for styling + parses overview text for topics"""

    def __init__(self):
        self.template_path = None
        self.presentation = None
        self.overview_topics = []
        self.theme_colors = {}
        self.fonts_config = {}
        self.presentation_title = ""
        self.presentation_subtitle = ""

    def load_and_analyze_template(self, template_path: str) -> Dict[str, Any]:
        """Load PPT template and extract styling"""
        try:
            if not os.path.exists(template_path):
                logger.error(f"Template file not found: {template_path}")
                return {}

            self.template_path = template_path
            self.presentation = Presentation(template_path)

            logger.info(f"✅ Template loaded: {template_path}")
            logger.info(f"📊 Total slides in template: {len(self.presentation.slides)}")

            styling = self._extract_template_styling()
            return styling

        except Exception as e:
            logger.error(f"❌ Error loading template: {str(e)}")
            return {}

    def _extract_template_styling(self) -> Dict[str, Any]:
        """Extract design elements (colors, fonts, layouts) from template"""
        try:
            styling = {
                "colors": {},
                "fonts": {},
                "layouts": [],
                "slide_width": self.presentation.slide_width,
                "slide_height": self.presentation.slide_height
            }

            for layout in self.presentation.slide_layouts:
                styling["layouts"].append({
                    "name": layout.name,
                    "slide_count": len([s for s in self.presentation.slides if s.slide_layout == layout])
                })

            slide_count = min(3, len(self.presentation.slides))
            for slide_idx in range(slide_count):
                slide = self.presentation.slides[slide_idx]
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.name:
                                    styling["fonts"][run.font.name] = True
                                try:
                                    if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                        color_hex = str(run.font.color.rgb)
                                        styling["colors"][color_hex] = True
                                except Exception:
                                    pass

            logger.info(f"✅ Template styling extracted")
            return styling

        except Exception as e:
            logger.error(f"❌ Error extracting styling: {str(e)}")
            return {}

    def parse_overview_text(self, overview_text: str) -> Dict[str, Any]:
        """
        Parse presentation overview from text input with AUTO-DETECTION
        Handles:
        - Explicit markers (TITLE:, TOPICS:)
        - Numbered/bulleted lists
        - Line-separated items (each on new line)
        - Space-separated items (fallback)
        """
        try:
            if not overview_text or len(overview_text.strip()) < 3:
                logger.error("❌ Overview text too short")
                return {}

            logger.info(f"📋 Parsing overview text...")

            overview_data = {
                "title": "",
                "subtitle": "",
                "topics": [],
                "raw_text": overview_text
            }

            # Try line-separated first
            lines = [line.strip() for line in overview_text.split('\n') if line.strip()]
            
            if not lines:
                logger.error("❌ No non-empty lines found")
                return {}

            logger.info(f"   Found {len(lines)} non-empty lines\n")

            title_found = False
            in_topics_section = False

            # STEP 1: Extract title and look for explicit sections
            for idx, line in enumerate(lines):
                line_clean = line.strip()
                if not line_clean:
                    continue
                # Explicit markers
                if line_clean.upper().startswith('TITLE:'):
                    overview_data["title"] = line_clean.split(':', 1)[1].strip()
                    title_found = True
                    logger.info(f"   ✓ Found TITLE marker: {overview_data['title']}")
                    continue
                if line_clean.upper().startswith('SUBTITLE:'):
                    overview_data["subtitle"] = line_clean.split(':', 1)[1].strip()
                    logger.info(f"   ✓ Found SUBTITLE marker: {overview_data['subtitle']}")
                    continue
                if line_clean.upper().startswith('TOPICS:') or line_clean.upper() == 'TOPICS':
                    in_topics_section = True
                    logger.info(f"   ✓ Found TOPICS section marker")
                    continue
                if not title_found and not overview_data["title"] and idx == 0:
                    if not self._is_marker_line(line_clean):
                        overview_data["title"] = line_clean
                        title_found = True
                        logger.info(f"   ✓ Using first line as TITLE: {overview_data['title']}")
                        continue

            # STEP 2: Extract topics (numbered/bulleted and in topics section)
            logger.info(f"\n   Extracting topics...")
            for line in lines:
                line_clean = line.strip()
                if not line_clean or self._is_marker_line(line_clean):
                    continue
                if self._is_topic_line(line_clean) or in_topics_section:
                    topic = self._clean_topic_line(line_clean)
                    if len(topic) > 2:
                        overview_data["topics"].append(topic)
                        logger.info(f"   ✓ Found TOPIC: {topic[:60]}...")

            # STEP 3: AUTO-DETECTION - fallback for plain lines (one per line)
            if not overview_data["topics"] and len(lines) > 1:
                logger.info(f"\n   ⚠️ No formatted topics found, AUTO-DETECTING (line-separated)...")
                for idx, line in enumerate(lines):
                    line_clean = line.strip()
                    if not line_clean or self._is_marker_line(line_clean):
                        continue
                    # Skip title line
                    if idx == 0 and overview_data["title"] == line_clean:
                        continue
                    # Skip subtitle line
                    if overview_data["subtitle"] and line_clean == overview_data["subtitle"]:
                        continue
                    if len(line_clean) > 2 and len(line_clean) < 200:
                        overview_data["topics"].append(line_clean)
                        logger.info(f"   ✓ AUTO-DETECTED TOPIC: {line_clean[:60]}...")

            # STEP 4: FALLBACK - Split by spaces if still no topics (all on one line)
            if not overview_data["topics"] and len(lines) == 1:
                logger.info(f"\n   ⚠️ All text on ONE line, splitting by spaces...")
                line_clean = lines[0]
                
                # Remove markers
                if line_clean.upper().startswith('TITLE:'):
                    overview_data["title"] = line_clean.split(':', 1)[1].strip()
                    line_clean = ""
                
                if line_clean:
                    # Split by spaces and filter out very short items
                    words = line_clean.split()
                    
                    # First item is likely title if not already set
                    if not overview_data["title"] and len(words) > 0:
                        # Take first 2-3 words as title if they look like a title
                        potential_title = words[0]
                        if len(potential_title) > 3:
                            overview_data["title"] = potential_title
                            words = words[1:]
                            logger.info(f"   ✓ Using first word as TITLE: {overview_data['title']}")
                    
                    # Rest are topics
                    for word in words:
                        if len(word) > 2 and len(word) < 50:
                            overview_data["topics"].append(word)
                            logger.info(f"   ✓ FALLBACK TOPIC: {word}")

            self.presentation_title = overview_data["title"]
            self.presentation_subtitle = overview_data["subtitle"]
            self.overview_topics = overview_data["topics"]

            logger.info(f"\n✅ Overview parsed:")
            logger.info(f"   Title: '{overview_data['title']}'")
            logger.info(f"   Subtitle: '{overview_data['subtitle']}'")
            logger.info(f"   Topics found: {len(overview_data['topics'])}")

            if not overview_data['topics']:
                logger.warning("\n⚠️ Still no topics found!")
                logger.warning(f"   Text lines: {lines}")
                return {}

            return overview_data

        except Exception as e:
            logger.error(f"❌ Error parsing overview: {str(e)}")
            import traceback
            traceback.print_exc()
            return {}

    def _is_marker_line(self, line: str) -> bool:
        """Check if line is a header/marker line"""
        markers = ['TITLE', 'SUBTITLE', 'TOPICS', 'OVERVIEW', 'PRESENTATION']
        return any(line.upper().startswith(m) for m in markers)

    def _is_topic_line(self, line: str) -> bool:
        """Check if line is a formatted topic (numbered, bulleted)"""
        # Numbered: 1. Topic, 2) Topic, 1: Topic
        if re.match(r'^[\d]+[\.\):]\s+', line):
            return True
        # Bulleted: - Topic, • Topic, * Topic, + Topic
        if re.match(r'^[-•*+]\s+', line):
            return True
        return False

    def _clean_topic_line(self, line: str) -> str:
        """Clean topic line (remove numbering, bullets)"""
        cleaned = re.sub(r'^[\d]+[\.\):]\s*', '', line).strip()
        cleaned = re.sub(r'^[-•*+]\s+', '', cleaned).strip()
        return cleaned

    def get_overview_topics(self) -> List[str]:
        """Get list of topics"""
        return self.overview_topics

    def analyze_full(self, template_path: str, overview_text: str) -> Dict[str, Any]:
        """Complete analysis: template + overview"""
        logger.info(f"\n{'='*60}")
        logger.info(f"📊 ANALYZING PPT TEMPLATE & OVERVIEW")
        logger.info(f"{'='*60}\n")

        # Load template
        logger.info("📋 Step 1: Loading template...")
        template_styling = self.load_and_analyze_template(template_path)

        if not template_styling:
            logger.error("❌ Template loading failed")
            return {}

        logger.info(f"✅ Template loaded\n")

        # Parse overview
        logger.info("📋 Step 2: Parsing overview text...")
        overview = self.parse_overview_text(overview_text)

        if not overview:
            logger.error("❌ Overview parsing failed")
            return {}

        if not overview.get("topics"):
            logger.error("❌ No topics found in overview")
            return {}

        logger.info(f"✅ Overview parsed\n")

        # Combine analysis
        analysis = {
            "analysis_complete": True,
            "template_path": template_path,
            "presentation_title": overview.get("title", "Untitled"),
            "presentation_subtitle": overview.get("subtitle", ""),
            "topics": overview.get("topics", []),
            "total_topics": len(overview.get("topics", [])),
            "template_styling": template_styling,
            "status": "✅ Complete analysis successful"
        }

        logger.info(f"{'='*60}")
        logger.info(f"📊 ANALYSIS COMPLETE")
        logger.info(f"{'='*60}")
        logger.info(f"Title: {analysis['presentation_title']}")
        logger.info(f"Topics: {analysis['total_topics']}")
        logger.info(f"{'='*60}\n")

        return analysis