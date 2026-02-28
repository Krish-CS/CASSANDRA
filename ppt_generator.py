"""
PPT Generator Module
Generates professional PowerPoint presentations using python-pptx
With image generation support via Pollinations AI
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from typing import Dict, List, Any, Tuple
import logging
import os
import io
import requests
import tempfile

# Optional: Image generator (may not be available)
try:
    from src.image_generator import ImageGenerator
    IMAGE_GENERATOR_AVAILABLE = True
except ImportError:
    IMAGE_GENERATOR_AVAILABLE = False
    ImageGenerator = None

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Generates professional PowerPoint presentations with image support"""
    
    def __init__(self):
        self.default_font = "Times New Roman"
        self.default_font_size = 20  # Content font size (reduced for better fit)
        self.title_font_size = 28    # Slide title font size
        self.heading_font_size = 28  # Slide title font size
        self.image_generator = ImageGenerator() if IMAGE_GENERATOR_AVAILABLE else None

        self.code_analysis = {}  # Will be set during generation
        self.project_name = ""  # Will be set during generation
        
        # Colors - ALL BLACK
        self.title_color = RGBColor(0, 0, 0)  # Black
        self.text_color = RGBColor(0, 0, 0)  # Black
        self.accent_color = RGBColor(0, 0, 0)  # Black
    
    def _set_shape_transparency(self, shape, transparency_percent):
        """Set transparency on a shape (0 = opaque, 100 = fully transparent)
        Uses OxmlElement to add alpha element to the shape's fill"""
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            
            # Get the spPr element
            spPr = shape._element.spPr
            solidFill = spPr.find(qn('a:solidFill'))
            
            if solidFill is not None:
                # Find srgbClr or schemeClr
                color_elem = solidFill.find(qn('a:srgbClr'))
                if color_elem is None:
                    color_elem = solidFill.find(qn('a:schemeClr'))
                
                if color_elem is not None:
                    # Remove existing alpha if present
                    existing_alpha = color_elem.find(qn('a:alpha'))
                    if existing_alpha is not None:
                        color_elem.remove(existing_alpha)
                    
                    # Add new alpha element (value is in 1000ths of a percent)
                    # 0% transparency = 100000, 25% transparency = 75000, etc.
                    alpha_value = int((100 - transparency_percent) * 1000)
                    alpha_elem = etree.SubElement(color_elem, qn('a:alpha'))
                    alpha_elem.set('val', str(alpha_value))
        except Exception as e:
            print(f"   ⚠️ Could not set transparency: {e}")
    
    def _download_background_image(self, url: str) -> str:
        """Download background image from URL and save to temp file"""
        try:
            print(f"   🌐 Downloading background image...")
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
            temp_file.write(response.content)
            temp_file.close()
            
            print(f"   ✅ Background image downloaded")
            return temp_file.name
        except Exception as e:
            print(f"   ⚠️ Could not download background: {e}")
            return None
    
    def _apply_background_to_slide(self, slide):
        """Apply background image to a slide"""
        if not self.background_image_path or not os.path.exists(self.background_image_path):
            return
        
        try:
            from pptx.oxml.ns import nsmap
            from lxml import etree
            
            # Get the slide's background
            # Add background shape as full-slide image
            left = Emu(0)
            top = Emu(0)
            width = self.slide_width
            height = self.slide_height
            
            # Add picture as background (insert at beginning so it's behind other shapes)
            pic = slide.shapes.add_picture(
                self.background_image_path,
                left, top, width, height
            )
            
            # Move picture to back
            spTree = slide.shapes._spTree
            sp = pic._element
            spTree.remove(sp)
            spTree.insert(2, sp)  # Insert after nvGrpSpPr and grpSpPr
            
        except Exception as e:
            print(f"   ⚠️ Could not apply background to slide: {e}")
    
    def generate_ppt(
        self,
        template_path: str,
        project_name: str,
        generated_content: Dict[str, Any],
        sections_config: Dict[str, Any],  # {sections: {name: style}, bullet_symbol: "➢"}
        output_path: str
    ) -> str:
        """
        Generate PowerPoint presentation using Cassandra template
        
        Args:
            template_path: Path to user's PPT template (.pptx) - used for backup
            project_name: Name of the project
            generated_content: AI-generated content
            sections_config: Dict with 'sections' mapping and 'bullet_symbol'
            output_path: Path to save the PPT
        
        Returns:
            Path to the generated PPT file
        """
        print("\n   🎨 Opening Cassandra PPT template...")
        
        # Extract config
        self.section_styles = sections_config.get('sections', {})
        self.bullet_symbol = sections_config.get('bullet_symbol', '➣')
        self.background_url = sections_config.get('background_url', '')
        self.thank_you_image_url = sections_config.get('thank_you_image_url', '')
        self.project_name = project_name
        self.code_analysis = generated_content.get('code_analysis', {})
        self.background_image_path = None
        
        # Download background image if provided
        if self.background_url:
            self.background_image_path = self._download_background_image(self.background_url)
        
        try:
            # Use Cassandra template as base
            cassandra_template = os.path.join(os.path.dirname(__file__), 'static', 'ppt', 'cassandra.pptx')
            
            # Fallback to user template if Cassandra template not found
            if os.path.exists(cassandra_template):
                prs = Presentation(cassandra_template)
                print("   ✅ Loaded Cassandra template")
            else:
                prs = Presentation(template_path)
                print(f"   ✅ Loaded user template (Cassandra template not found)")
            
            # Get slide dimensions from template
            self.slide_width = prs.slide_width
            self.slide_height = prs.slide_height
            
            # Delete the first placeholder slide if exists
            if len(prs.slides) > 0:
                slide_id = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(slide_id)
                del prs.slides._sldIdLst[0]
                print("   🗑️ Deleted placeholder slide")
            
            # Add generated slides - EXACTLY as shown in preview

            print("   📝 Adding generated slides...")
            
            # Loop through ALL chapters from preview - add them AS-IS
            for chapter in generated_content.get("chapters", []):
                chapter_title = chapter.get("title", "")
                
                for section in chapter.get("sections", []):
                    section_title = section.get("title", "")
                    content = section.get("content", "")
                    
                    if not content or len(content.strip()) < 20:
                        continue
                    
                    # Use section style if provided, otherwise infer from title
                    section_style = section.get("style", None)
                    if section_style:
                        style = section_style
                    else:
                        match_title = section_title if section_title else chapter_title
                        style, _, _, _ = self._get_config_for_section(match_title)
                    
                    # Create slide title - use the chapter title
                    slide_title = chapter_title.upper()
                    
                    # Add content slide (skip if style is 'none')
                    if style != 'none':
                        self._add_content_slide(prs, slide_title, content, style=style)
                        print(f"   ✅ Slide added: {slide_title[:50]}...")
                    else:
                        print(f"   ⏭️ Skipped content for: {slide_title} (style=none)")

            
            # Add Thank You Slide (extra slide at the end)
            self._add_thank_you_slide(prs)
            print("   ✅ Thank you slide added")

            
            # Count total slides
            total_slides = len(prs.slides)
            print(f"\n   📊 Generated {total_slides} slides")

            
            # Save presentation
            prs.save(output_path)
            
            if os.path.exists(output_path):
                file_size = os.path.getsize(output_path)
                print(f"   ✅ PPT saved successfully ({file_size:,} bytes)")
            
            return output_path
            
        except Exception as e:
            print(f"   ❌ Error generating PPT: {str(e)}")
            import traceback
            traceback.print_exc()
            raise
    
    def _add_title_slide(self, prs: Presentation, project_name: str):
        """Add title slide"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title text box
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(12.333)
        height = Inches(2)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = project_name.upper()
        p.font.name = self.default_font
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.title_color
        p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        p2 = tf.add_paragraph()
        p2.text = "Project Presentation"
        p2.font.name = self.default_font
        p2.font.size = Pt(24)
        p2.font.color.rgb = self.accent_color
        p2.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(
        self, 
        prs: Presentation, 
        title: str, 
        content: str, 
        style: str = "bullet"
    ):
        """Add content slide with bullet or paragraph style"""
        # Use layout 10 (BLANK) from Cassandra template
        try:
            slide_layout = prs.slide_layouts[10]  # BLANK layout
        except:
            slide_layout = prs.slide_layouts[6]
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Apply background image first (so it's behind content)
        self._apply_background_to_slide(slide)
        
        # Remove any placeholder shapes (like "Click to add title")
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)
        
        # ==================
        # TITLE BOX (Top)
        # ==================
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12.33)
        title_height = Inches(0.7)
        
        # Title background - subtle rounded corners (matches preview)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            title_left, title_top, title_width, title_height
        )
        # Set subtle corner radius for title
        try:
            title_bg.adjustments[0] = 0.1  # Subtle corners for smaller box
        except:
            pass
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Only add transparency if background image exists
        if self.background_image_path:
            self._set_shape_transparency(title_bg, 20)
        title_bg.line.fill.background()

        
        # Title text
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        title_p = title_tf.paragraphs[0]
        title_p.text = title.upper()
        title_p.font.name = self.default_font
        title_p.font.size = Pt(self.heading_font_size)
        title_p.font.bold = True
        title_p.font.color.rgb = self.text_color
        title_p.alignment = PP_ALIGN.CENTER
        
        # ==================
        # CONTENT BOX (Below title with gap)
        # ==================
        content_left = Inches(0.5)
        content_top = Inches(1.2)  # Gap after title (0.3 + 0.7 + 0.2 gap = 1.2)
        content_width = Inches(12.33)
        content_height = Inches(5.8)  # Rest of slide height
        
        # Content background - subtle rounded corners (matches preview)
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            content_left, content_top, content_width, content_height
        )
        # Set subtle corner radius (10% of smaller dimension)
        try:
            content_bg.adjustments[0] = 0.02  # Very subtle corners
        except:
            pass
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Only add transparency if background image exists
        if self.background_image_path:
            self._set_shape_transparency(content_bg, 20)
        content_bg.line.fill.background()

        
        # Content text
        content_box = slide.shapes.add_textbox(
            content_left + Inches(0.2),
            content_top + Inches(0.15),
            content_width - Inches(0.4),
            content_height - Inches(0.3)
        )
        content_tf = content_box.text_frame
        content_tf.word_wrap = True
        content_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        content_tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center content vertically
        
        content_tf.margin_left = Inches(0.1)
        content_tf.margin_right = Inches(0.1)
        content_tf.margin_top = Inches(0.1)
        content_tf.margin_bottom = Inches(0.1)
        
        # ==================
        # ADD TEXT CONTENT
        # ==================
        if style == "bullet":
            points = self._extract_bullet_points(content)
            count = len(points)
            
            # Fixed font size 20pt for all bullet points
            font_size = 20
            
            for i, point in enumerate(points):
                para = content_tf.paragraphs[0] if i == 0 else content_tf.add_paragraph()
                # Use selected bullet symbol
                para.text = f"{self.bullet_symbol} {point}"

                para.font.name = self.default_font
                para.font.size = Pt(font_size)
                para.font.color.rgb = self.text_color
                para.alignment = PP_ALIGN.LEFT
                para.line_spacing = 1.3
                para.space_before = Pt(6)
                para.space_after = Pt(6)
        else:

            p = content_tf.paragraphs[0]
            p.text = self._clean_for_slide(content)
            p.font.name = self.default_font
            p.font.size = Pt(self.default_font_size)
            p.font.color.rgb = self.text_color
            p.alignment = PP_ALIGN.JUSTIFY
            p.line_spacing = 1.3
    
    def _enable_autofit(self, textbox):
        """Enable PowerPoint's auto-fit (shrink text on overflow) via XML"""
        try:
            from pptx.oxml.ns import qn
            from lxml import etree
            
            # Get the textBody (txBody) element
            txBody = textbox._element.find(qn('p:txBody'))
            if txBody is None:
                return
            
            # Get or create bodyPr
            bodyPr = txBody.find(qn('a:bodyPr'))
            if bodyPr is None:
                return
            
            # Remove existing fit settings
            for child in list(bodyPr):
                tag = child.tag.split('}')[-1]
                if tag in ['normAutofit', 'spAutoFit', 'noAutofit']:
                    bodyPr.remove(child)
            
            # Add normAutofit (shrink text to fit)
            normAutofit = etree.SubElement(bodyPr, qn('a:normAutofit'))
            normAutofit.set('fontScale', '70000')  # Allow shrinking to 70%
            normAutofit.set('lnSpcReduction', '20000')  # Allow 20% line spacing reduction
            
        except Exception as e:
            print(f"   ⚠️ Could not enable autofit: {e}")
    
    def _add_bullet_content(self, text_frame, content: str, num_bullets: int = 6):
        """Add content as bullet points with dynamic font sizing"""
        # Extract bullet points from content
        points = self._extract_bullet_points(content)
        
        # FIX 4: Scale font size by number of bullets
        # More bullets = smaller font (18-22 range)
        if num_bullets <= 4:
            font_size = 22
            line_spacing = 1.4
            space_between = 8
        elif num_bullets <= 6:
            font_size = 20
            line_spacing = 1.3
            space_between = 6
        else:
            font_size = 18
            line_spacing = 1.2
            space_between = 4

        
        for i, point in enumerate(points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # Add bullet symbol (use configured symbol)
            p.text = f"{self.bullet_symbol}  {point}"
            p.font.name = self.default_font
            p.font.size = Pt(font_size)  # Dynamic font size
            p.font.color.rgb = self.text_color
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing  # Dynamic line spacing
            p.space_before = Pt(space_between)  # Dynamic spacing
            p.space_after = Pt(space_between)

    
    def _add_paragraph_content(self, text_frame, content: str):
        """Add content as paragraph"""
        # Clean content and make it concise for slides
        clean_content = self._clean_for_slide(content)
        
        p = text_frame.paragraphs[0]
        p.text = clean_content
        p.font.name = self.default_font
        p.font.size = Pt(self.default_font_size)
        p.font.color.rgb = self.text_color
        p.alignment = PP_ALIGN.JUSTIFY
        p.line_spacing = 1.3  # Reduced from 1.5 for better fit

    
    def _fetch_thank_you_image(self) -> str:
        """Fetch a random thank you image from Pexels - tries multiple search terms"""
        import random
        
        # Try multiple search terms for variety
        search_terms = ['thank you', 'gratitude', 'appreciation', 'colorful abstract', 'beautiful nature']
        random.shuffle(search_terms)
        
        try:
            pexels_api_key = os.getenv('PEXELS_API_KEY', 'mVCnJKoyP7wAJQIJ3cPgaIpDwxdnZqjETR3gR2qPwySdVLV0D4VnzPxk')
            headers = {'Authorization': pexels_api_key}
            
            for search_term in search_terms[:3]:  # Try up to 3 terms
                print(f"   🔍 Searching Pexels for: {search_term}")
                response = requests.get(
                    'https://api.pexels.com/v1/search',
                    params={'query': search_term, 'per_page': 20, 'orientation': 'landscape'},
                    headers=headers,
                    timeout=15
                )
                
                if response.status_code == 200:
                    data = response.json()
                    photos = data.get('photos', [])
                    if photos:
                        # Get random photo
                        photo = random.choice(photos)
                        image_url = photo.get('src', {}).get('large2x') or photo.get('src', {}).get('original')
                        
                        if image_url:
                            # Download the image
                            img_response = requests.get(image_url, timeout=30)
                            if img_response.status_code == 200:
                                temp_file = tempfile.NamedTemporaryFile(suffix='.jpg', delete=False)
                                temp_file.write(img_response.content)
                                temp_file.close()
                                print(f"   ✅ Thank You image fetched from Pexels ({search_term})")
                                return temp_file.name
        except Exception as e:
            print(f"   ⚠️ Could not fetch Thank You image: {e}")
        
        return None

    
    def _add_thank_you_slide(self, prs: Presentation):
        """Add thank you slide with background image and centered text"""
        try:
            slide_layout = prs.slide_layouts[6]  # Blank layout
        except:
            slide_layout = prs.slide_layouts[0]
        
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove any placeholders to ensure clean slate
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                sp = shape._element
                sp.getparent().remove(sp)
        
        # 1. Add Background
        thank_you_image = None
        
        # Try custom Image URL first
        if self.thank_you_image_url:
            thank_you_image = self._download_background_image(self.thank_you_image_url)
            
        # If no custom image and NOT using a template, try random Pexels
        if not thank_you_image and not self.background_image_path:
            thank_you_image = self._fetch_thank_you_image()
            
        if thank_you_image and os.path.exists(thank_you_image):
            # Add full-slide background image - NO TEXT
            pic = slide.shapes.add_picture(
                thank_you_image,
                Emu(0), Emu(0),
                self.slide_width, self.slide_height
            )
            
            # Send to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            
            # Cleanup
            try:
                os.unlink(thank_you_image)
            except: pass
        else:
            # Fallback to template background + Centered Text
            self._apply_background_to_slide(slide)
            
            # Add "THANK YOU" Text Overlay (IMPROVED CENTERING)
            # Use a centered semi-transparent white box
            box_width = Inches(10)  # Moderate width to look "centered"
            box_height = Inches(2.5)
            
            left = (self.slide_width - box_width) / 2
            top = (self.slide_height - box_height) / 2
            
            # Add background shape for text
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, box_width, box_height
            )
            try:
                shape.adjustments[0] = 0.05
            except: pass
            
            # Style the box: White with transparency
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            self._set_shape_transparency(shape, 20)  # 20% transparent
            shape.line.fill.background()  # No border
            
            # Add Text with proper alignment
            tf = shape.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically Center
            
            p = tf.paragraphs[0]
            p.text = "THANK YOU"
            p.font.name = self.default_font
            p.font.size = Pt(66)  # Large, bold
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Always black on white box
            p.alignment = PP_ALIGN.CENTER  # Horizontally Center


    
    def _extract_bullet_points(self, content: str) -> List[str]:
        """Extract or create bullet points from content - complete sentences only"""
        # First try to split by newlines
        lines = content.split('\n')
        points = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Remove ALL existing bullet markers (comprehensive list including preview options)
            # This includes: ✱ ✲ ❖ ✦ ✧ ✤ and all other common bullets
            line = line.lstrip('•➢➣▪▫-*►○●⁃◆◇■□▸▹▶▷→➤✓✔★☆◉⇒❥☸✦✧⊳⊲⫸⫷⪢⪡·⊛◌◍◎◘◦☉⁌⁍◈☐☑☒❧☙✤✱✲❖↠↣↦↬⇛⇝⇢⇨➙➛➜➝➞➟➠➡➥➦➧➨➮➱➲➳➵➸➼➽➾⇾‣▻ ')

            line = line.lstrip('0123456789.')
            line = line.strip()
            
            if line and len(line) > 10:
                points.append(line)
        
        # If no clear points from lines, split by sentences
        if len(points) < 2:
            # Split content into sentences
            content_clean = content.replace('\n', ' ')
            sentences = []
            current = ""
            
            for char in content_clean:
                current += char
                if char in '.!?':
                    sent = current.strip()
                    if sent and len(sent) > 20:
                        sentences.append(sent)
                    current = ""
            
            points = sentences if sentences else [content_clean]
        
        # Clean up points - ensure they end properly
        clean_points = []
        for point in points:
            point = point.strip()
            if not point:
                continue
            
            # Limit length - find last complete word/sentence within limit
            if len(point) > 120:
                # Find last period, comma, or space before limit
                cut_point = 120
                for i in range(120, 50, -1):
                    if point[i] in '.!?':
                        cut_point = i + 1
                        break
                    elif point[i] == ' ':
                        cut_point = i
                        break
                point = point[:cut_point].strip()
                # Ensure it ends with proper punctuation
                if point and point[-1] not in '.!?':
                    point = point + "."
            
            # Ensure ends with punctuation
            if point and point[-1] not in '.!?':
                point = point + "."
            
            clean_points.append(point)
        
        # Limit to 8 points max per slide
        return clean_points[:8]

    
    def _clean_for_slide(self, content: str) -> str:
        """Clean and shorten content for slide display - complete sentences only"""
        # Remove extra whitespace
        content = ' '.join(content.split())
        
        # If content is short enough, just ensure it ends properly
        if len(content) <= 1000:
            if content and content[-1] not in '.!?':
                content = content + "."
            return content
        
        # For longer content, find a good cut point at a sentence boundary
        cut_content = content[:1000]
        
        # Find last sentence ending
        last_period = -1
        for i in range(len(cut_content) - 1, 200, -1):
            if cut_content[i] in '.!?':
                last_period = i + 1
                break
        
        if last_period > 200:
            return cut_content[:last_period].strip()
        else:
            # No good sentence break found, cut at last space and add period
            last_space = cut_content.rfind(' ', 200, 1000)
            if last_space > 200:
                result = cut_content[:last_space].strip()
                if result[-1] not in '.!?':
                    result = result + "."
                return result
            else:
                return cut_content.strip() + "."

    
    def _determine_style(
        self, 
        section_title: str, 
        sections_config: Dict[str, str],
        content: str
    ) -> str:
        """Determine if section should be bullet or paragraph"""
        result = self._get_config_for_section(section_title)
        return result[0]
    
    def _get_config_for_section(self, section_title: str) -> Tuple[str, bool, bool, str, str]:
        """Get style and image config for a section from config
        
        Returns:
            Tuple of (style, ai_image_enabled, custom_image_enabled, custom_image_data, custom_caption)
        
        Returns:
            Tuple of (style, image_enabled)
        """
        section_lower = section_title.lower()
        
        # Check config first - handle both string and dict formats
        for key, config_value in self.section_styles.items():
            key_lower = key.lower()
            # Match by removing section numbers
            key_words = [w for w in key_lower.split() if not w.replace('.', '').isdigit()]
            section_words = [w for w in section_lower.split() if not w.replace('.', '').isdigit()]
            key_clean = ' '.join(key_words)
            section_clean = ' '.join(section_words)
            
            if (key_lower in section_lower or section_lower in key_lower or
                key_clean == section_clean or 
                (key_clean and section_clean and (key_clean in section_clean or section_clean in key_clean))):
                
                if isinstance(config_value, dict):
                    style = config_value.get('style', 'paragraph')
                    ai_image = config_value.get('image', False)
                    custom_image = config_value.get('customImage', False)
                    custom_images_data = config_value.get('customImagesData', [])  # Array format
                    return style, ai_image, custom_image, custom_images_data
                else:
                    return config_value, False, False, []
        
        # Default logic based on section type - most slides should be bullet
        paragraph_keywords = [
            "abstract", "introduction", "conclusion", "summary",
            "overview", "background", "description"
        ]
        
        for keyword in paragraph_keywords:
            if keyword in section_lower:
                return "paragraph", False, False, []
        
        # Default to bullet style for all other slides
        return "bullet", False, False, []
    
    def _get_style_for_section(self, section_title: str) -> str:
        """Get style for a section from config (backwards compat)"""
        result = self._get_config_for_section(section_title)
        return result[0]  # Return just the style
    
    def _add_image_slide(self, prs: Presentation, title: str, content: str):
        """Add an image slide with generated diagram"""
        print(f"    🖼️ Generating image slide for: {title}")
        
        try:
            # Generate image using Pollinations AI
            image_bytes = self.image_generator.generate_section_image_sync(
                section_title=title,
                section_content=content,
                code_analysis=self.code_analysis or {},
                project_name=self.project_name or "Project"
            )
            
            if not image_bytes:
                print(f"    ⚠️ No image generated for {title}")
                return
            
            # Create the slide
            try:
                slide_layout = prs.slide_layouts[6]  # Blank layout
            except:
                slide_layout = prs.slide_layouts[0]
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove any placeholder shapes
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.is_placeholder:
                    shapes_to_remove.append(shape)
            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
            
            # Add title
            title_left = Inches(0.5)
            title_top = Inches(0.3)
            title_width = Inches(12.333)
            title_height = Inches(0.8)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_p = title_tf.paragraphs[0]
            title_p.text = title
            title_p.font.name = self.default_font
            title_p.font.size = Pt(self.heading_font_size)
            title_p.font.bold = True
            title_p.font.color.rgb = self.text_color
            title_p.alignment = PP_ALIGN.CENTER
            
            # Add image - centered on slide, properly sized to fit
            # Standard slide is 13.333 x 7.5 inches, so image should be smaller
            image_stream = io.BytesIO(image_bytes)
            image_width = Inches(8.0)  # Reduced from 11.333 to fit slide
            image_height = Inches(4.5)  # Proportional height
            image_left = Inches(2.67)  # Center: (13.333 - 8) / 2 = 2.67
            image_top = Inches(1.3)
            
            slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width, height=image_height)
            
            # Add caption below image
            caption_top = Inches(6.0)  # Adjusted to be below image
            caption_box = slide.shapes.add_textbox(Inches(0.5), caption_top, Inches(12.333), Inches(0.5))
            caption_tf = caption_box.text_frame
            caption_p = caption_tf.paragraphs[0]
            caption_p.text = f"{title} - Diagram"
            caption_p.font.name = self.default_font
            caption_p.font.size = Pt(14)
            caption_p.font.bold = True
            caption_p.font.color.rgb = self.text_color
            caption_p.alignment = PP_ALIGN.CENTER
            
            print(f"    ✅ Image slide added: {title}")
            
        except Exception as e:
            print(f"    ⚠️ Error adding image slide: {e}")
            import traceback
            traceback.print_exc()
    
    def _add_custom_image_slide(self, prs: Presentation, title: str, image_data: str, custom_caption: str = ""):
        """Add an image slide with custom uploaded image"""
        print(f"    🖼️ Inserting custom image slide for: {title}")
        
        try:
            import base64
            
            # Decode base64 image (remove data URL prefix if present)
            if ',' in image_data:
                image_data = image_data.split(',')[1]
            image_bytes = base64.b64decode(image_data)
            
            # Create the slide
            try:
                slide_layout = prs.slide_layouts[6]  # Blank layout
            except:
                slide_layout = prs.slide_layouts[0]
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove any placeholder shapes
            shapes_to_remove = []
            for shape in slide.shapes:
                if shape.is_placeholder:
                    shapes_to_remove.append(shape)
            for shape in shapes_to_remove:
                sp = shape._element
                sp.getparent().remove(sp)
            
            # Add title
            title_left = Inches(0.5)
            title_top = Inches(0.3)
            title_width = Inches(12.333)
            title_height = Inches(0.8)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_tf = title_box.text_frame
            title_tf.word_wrap = True
            title_p = title_tf.paragraphs[0]
            title_p.text = title
            title_p.font.name = self.default_font
            title_p.font.size = Pt(self.heading_font_size)
            title_p.font.bold = True
            title_p.font.color.rgb = self.text_color
            title_p.alignment = PP_ALIGN.CENTER
            
            # Add image - centered on slide, properly sized to fit
            image_stream = io.BytesIO(image_bytes)
            image_width = Inches(8.0)
            image_height = Inches(4.5)
            image_left = Inches(2.67)
            image_top = Inches(1.3)
            
            slide.shapes.add_picture(image_stream, image_left, image_top, width=image_width, height=image_height)
            
            # Add caption below image
            caption_text = custom_caption.strip() if custom_caption.strip() else f"{title} - Diagram"
            caption_top = Inches(6.0)
            caption_box = slide.shapes.add_textbox(Inches(0.5), caption_top, Inches(12.333), Inches(0.5))
            caption_tf = caption_box.text_frame
            caption_p = caption_tf.paragraphs[0]
            caption_p.text = caption_text
            caption_p.font.name = self.default_font
            caption_p.font.size = Pt(14)
            caption_p.font.bold = True
            caption_p.font.color.rgb = self.text_color
            caption_p.alignment = PP_ALIGN.CENTER
            
            print(f"    ✅ Custom image slide added: {title}")
            
        except Exception as e:
            print(f"    ⚠️ Error adding custom image slide: {e}")
            import traceback
            traceback.print_exc()
    
    def _extract_objectives(self, generated_content: Dict[str, Any]) -> str:
        """Extract objectives from generated content"""
        for chapter in generated_content.get("chapters", []):
            for section in chapter.get("sections", []):
                if "objective" in section.get("title", "").lower():
                    return section.get("content", "")
        return ""
